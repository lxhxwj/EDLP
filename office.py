#!/usr/bin/python
# -*- coding: UTF-8 -*-
#载入必要的模块
import pygame
from docx import Document
from docx.shared import Inches
from multiprocessing import Pool
from pptx import Presentation
from pptx.util import Inches
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Cm
from pptx.enum.chart import XL_LEGEND_POSITION
from reportlab.pdfgen import canvas
from reportlab.platypus.tables import Table, TableStyle
from reportlab.lib.units import inch
from reportlab.platypus import Paragraph,Frame
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import Image as platImage
from PIL import Image
from reportlab.lib import colors
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfbase import pdfmetrics
import time
import random
import socket
import getopt
import sys
import os
import xlwt
import xlrd
import string
import struct
import _winreg

WLFS = 1

class office:

    def hello(self):
        print('hello')
    
    def CreateExcel(self,text,s):
        try:
            i = 0
            z = 3
            style = xlwt.XFStyle()
            font = xlwt.Font()
            font.name = 'Times New Roman'
            font.bold = True
            style.font = font
            filename = xlwt.Workbook(encoding = 'UTF-8')
            while i < s:
                t = 0
                x = random.randint(0, 10)
                y = random.randint(0, 10)
                sheet = filename.add_sheet("test %s" % i)
                while t < 100:
                    sheet.write(x,y,text,style)
                    x = x + 1
                    t = t + 1
                i = i + 1
            sheet = filename.add_sheet("test %s" % s)
            while z < 10:
                salt = ''.join(random.sample(string.ascii_letters + string.digits, 16))
                sheet.write(2,z - 3,"王宝强" + salt,style)
                z = z + 1
            filename.save("Result\\dlp.xls")
            return True
        except Exception,e:
            print(str(e))
            return False
    def CreateExcelNum(self,text,s,num):
        try:
            i = 0
            z = 3
            style = xlwt.XFStyle()
            font = xlwt.Font()
            font.name = 'Times New Roman'
            font.bold = True
            style.font = font
            filename = xlwt.Workbook(encoding = 'UTF-8')
            while i < s:
                t = 0
                x = random.randint(0, 10)
                y = random.randint(0, 10)
                sheet = filename.add_sheet("test %s" % i)
                while t < 100:
                    sheet.write(x,y,text,style)
                    x = x + 1
                    t = t + 1
                i = i + 1
            sheet = filename.add_sheet("test %s" % s)
            while z < 10:
                salt = ''.join(random.sample(string.ascii_letters + string.digits, 16))
                sheet.write(2,z - 3,"王宝强" + salt,style)
                z = z + 1
            filename.save("Result\\dlp%s.xls" % num)
            return True
        except Exception,e:
            print(str(e))
            return False
    def CreateExcelE(self,text,s,n):
        m = 0
        o = office()
        while m < n:
            o.CreateExcelNum(text,s,m)
            m = m + 1
        return True

    def CreateWord(self,s):
        i = 0
        z = 0
        document = Document()

        while i < s:
            document.add_heading(u'文档标题', 0)
            p = document.add_paragraph(u'这是一个自然段 ')
            p.add_run('bold').bold = True
            p.add_run(u' 还有 ')
            p.add_run('italic.').italic = True

            document.add_heading(u'1级别标题', level=1)
            document.add_paragraph(u'引用', style='IntenseQuote')

            document.add_paragraph(u'符号列表', style='ListBullet')
            document.add_paragraph(u'数字列表t', style='ListNumber')
            document.add_paragraph(u'我的微信:')
            document.add_picture('file\\tupian.jpg', width=Inches(3.25))

            table = document.add_table(rows=3, cols=3)
            hdr_cells = table.rows[0].cells
            hdr_cells[0].text = u'第一列'
            hdr_cells[1].text = u'第二列'
            hdr_cells[2].text = u'第三列'

            hdr_cells = table.rows[1].cells
            hdr_cells[0].text = '1'
            hdr_cells[1].text = '21'
            hdr_cells[2].text = 'qwertyuiop'

            hdr_cells = table.rows[2].cells
            hdr_cells[0].text = '2'
            hdr_cells[1].text = '43'
            hdr_cells[2].text = 'asdfghjkl'
            i = i + 1
       
        document.add_heading(u'我是关键字', 0)
        while z < 10:
            salt = ''.join(random.sample(string.ascii_letters + string.digits, 16))
            document.add_paragraph(u'王宝强' + salt)
            z = z + 1

        document.add_page_break()

        document.save("Result\\dlp.docx")

        return True
    def CreateWordNum(self,s,num):
        i = 0
        z = 0
        document = Document()

        while i < s:
            document.add_heading(u'文档标题', 0)
            p = document.add_paragraph(u'这是一个自然段 ')
            p.add_run('bold').bold = True
            p.add_run(u' 还有 ')
            p.add_run('italic.').italic = True

            document.add_heading(u'1级别标题', level=1)
            document.add_paragraph(u'引用', style='IntenseQuote')

            document.add_paragraph(u'符号列表', style='ListBullet')
            document.add_paragraph(u'数字列表t', style='ListNumber')
            document.add_paragraph(u'我的微信:')
            document.add_picture('file\\tupian.jpg', width=Inches(3.25))

            table = document.add_table(rows=3, cols=3)
            hdr_cells = table.rows[0].cells
            hdr_cells[0].text = u'第一列'
            hdr_cells[1].text = u'第二列'
            hdr_cells[2].text = u'第三列'

            hdr_cells = table.rows[1].cells
            hdr_cells[0].text = '1'
            hdr_cells[1].text = '21'
            hdr_cells[2].text = 'qwertyuiop'

            hdr_cells = table.rows[2].cells
            hdr_cells[0].text = '2'
            hdr_cells[1].text = '43'
            hdr_cells[2].text = 'asdfghjkl'
            i = i + 1
       
        document.add_heading(u'我是关键字', 0)
        while z < 10:
            salt = ''.join(random.sample(string.ascii_letters + string.digits, 16))
            document.add_paragraph(u'王宝强' + salt)
            z = z + 1

        document.add_page_break()

        document.save("Result\\dlp%s.docx" % num)

        return True
    def CreateWordE(self,s,n):
        m = 0
        o = office()
        while m < n:
            o.CreateWordNum(s,m)
            m = m + 1
        return True

    def CreatePpt(self,s):
        i = 0
        z = 0
        # 创建幻灯片 ------
        prs = Presentation()

        while i < s:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            # 定义图表数据 ---------------------
            chart_data = ChartData()
            chart_data.categories = ['East', 'West', 'Midwest']
            chart_data.add_series('Series 1', (19.2, 21.4, 16.7))
         
            # 将图表添加到幻灯片 --------------------
            x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
            slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
            )
            i = i + 1

        title_slide_layout1 = prs.slide_layouts[0]
        slide1 = prs.slides.add_slide(title_slide_layout1)
         
        title1 = slide1.shapes.title
        subtitle1 = slide1.placeholders[1]
         
        # 设置标题和副标题
        title1.text = "我是关键字"
        salt = ''.join(random.sample(string.ascii_letters + string.digits, 16))
        subtitle1.text = "王宝强" + salt

        prs.save('Result\\dlp.pptx')

        return True

    def CreatePptNum(self,s,num):
        i = 0
        z = 0
        # 创建幻灯片 ------
        prs = Presentation()

        while i < s:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            # 定义图表数据 ---------------------
            chart_data = ChartData()
            chart_data.categories = ['East', 'West', 'Midwest']
            chart_data.add_series('Series 1', (19.2, 21.4, 16.7))
         
            # 将图表添加到幻灯片 --------------------
            x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
            slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
            )
            i = i + 1

        title_slide_layout1 = prs.slide_layouts[0]
        slide1 = prs.slides.add_slide(title_slide_layout1)
         
        title1 = slide1.shapes.title
        subtitle1 = slide1.placeholders[1]
         
        # 设置标题和副标题
        title1.text = "我是关键字"
        salt = ''.join(random.sample(string.ascii_letters + string.digits, 16))
        subtitle1.text = "王宝强" + salt

        prs.save('Result\\dlp%s.pptx' % num)

        return True

    def CreatePptE(self,s,n):
        m = 0
        o = office()
        while m < n:
            o.CreatePptNum(s,m)
            m = m + 1
        return True

    def CreateTxt(self,s):
        first = []
        second = []
        f = open('Result\\dlp.txt','w')
        with open('file\\ftp.txt', 'r') as f1:
            for line in f1:
                line = line.strip()
                first.append(line)
            f1.close()
        with open('file\\ftp.txt', 'r') as f2:
            for line2 in f2:
                line2 = line2.strip()
                second.append(line2)
            f2.close()
        for i in range(0,s):
            salt = ''.join(random.sample(string.ascii_letters + string.digits, 16))
            result = first[i] + '\t' + second[i] + '\n'
            f.write(result)
        f.write(salt)
        f.close()
        return True

    def CreateTxtE(self,s,n):
        first = []
        second = []
        x = 0
        while x < n:
            f = open('Result\\dlp%s.txt' % x,'w')
            with open('file\\ftp.txt', 'r') as f1:
                for line in f1:
                    line = line.strip()
                    first.append(line)
                f1.close()
            with open('file\\ftp.txt', 'r') as f2:
                for line2 in f2:
                    line2 = line2.strip()
                    second.append(line2)
                f2.close()
            for i in range(0,s):
                salt = ''.join(random.sample(string.ascii_letters + string.digits, 16))
                result = first[i] + '\t' + second[i] + '\n'
                f.write(result)
            f.write(salt)
            f.close()
            x = x + 1
        return True

    def CreateImage(self):
        pygame.init()
        salt = ''.join(random.sample(string.ascii_letters + string.digits, 16))
        #待转换文字
        text = u"woshi啊实打实大家卡萨丁比较卡失败的健康萨比\
        大数据科技部科技表示大空间的巴萨健康大鼠标接口\
        12372193721983198diasdhiuadslkandklasnbdklasd王宝强" + salt
        #设置字体和字号
        font = pygame.font.SysFont('Microsoft YaHei', 64)
        #渲染图片，设置背景颜色和字体样式,前面的颜色是字体颜色
        ftext = font.render(text, True, (65, 83, 130),(255, 255, 255))
        #保存图片
        pygame.image.save(ftext,"Result\\dlp.png")#图片保存地址
        return True
    def CreateImageNum(self,num):
        pygame.init()
        salt = ''.join(random.sample(string.ascii_letters + string.digits, 16))
        #待转换文字
        text = u"woshi啊实打实大家卡萨丁比较卡失败的健康萨比\
        大数据科技部科技表示大空间的巴萨健康大鼠标接口\
        12372193721983198diasdhiuadslkandklasnbdklasd王宝强" + salt
        #设置字体和字号
        font = pygame.font.SysFont('Microsoft YaHei', 64)
        #渲染图片，设置背景颜色和字体样式,前面的颜色是字体颜色
        ftext = font.render(text, True, (65, 83, 130),(255, 255, 255))
        #保存图片
        pygame.image.save(ftext,"Result\\dlp%s.png" % num)#图片保存地址
        return True
    def CreateImageE(self,n):
        m = 0
        o = office()
        while m < n:
            o.CreateImageNum(m)
            m = m + 1
        return True

    def CreatePdf(self,s):
        x = 0
        salt = ''.join(random.sample(string.ascii_letters + string.digits, 16))
        #pdfmetrics.registerFont(TTFont('hei', 'hei.TTF'))
        #import testSubFun
        #testSubFun.testSubFunc('first')
        #设置页面大小
        c = canvas.Canvas('Result\\dlp.pdf',pagesize=A4)
        xlength,ylength = A4
        #print('width:%d high:%d'%(xlength,ylength))
        #c.line(1,1,ylength/2,ylength)
        #设置文字类型及字号
        #c.setFont('hei',20)
        #生成一个table表格
        atable = [[1,2,3,4,5,6,7,8],[11,12,13,14,15,16,17,18]]
        t = Table(atable,50,20)
        t.setStyle(TableStyle([('ALIGN',(0,0),(3,1),'CENTER'),
                               ('INNERGRID',(0,0),(-1,-1),0.25,colors.black),
                               ('BOX',(0,0),(-1,-1),0.25,colors.black)]))
        textOb = c.beginText(1,ylength-10)
        indexVlaue = 0
        while(indexVlaue < ylength):
            textStr = '''wo shi tupian---wo shi tupian--wo shi tupian--wo shi tupian%d'''%indexVlaue + salt
            #print('nextline,nextline%d'%indexVlaue)
            textOb.textLine(textStr)
            indexVlaue = indexVlaue + 1
            break
        c.drawText(textOb)
        
        #简单的图片载入
        imageValue = 'file\\dlp.png'
        c.drawImage(imageValue,97,97,650,650)
        #c.drawImage('file\\dlp.png',50,50,50,50)
        t.split(0,0)
        t.drawOn(c,100,1)
        c.showPage()
        #换页的方式不同的showPage
        while x < s:
            imageValue = 'file\\dlp.png'
            c.drawImage(imageValue,97,97,650,650)
            c.drawString(0,0,'tupian%s' % x)
            c.showPage()
            x = x + 1
        c.save()
    def CreatePdfNum(self,s,num):
        x = 0
        salt = ''.join(random.sample(string.ascii_letters + string.digits, 16))
        #pdfmetrics.registerFont(TTFont('hei', 'hei.TTF'))
        #import testSubFun
        #testSubFun.testSubFunc('first')
        #设置页面大小
        c = canvas.Canvas('Result\\dlp%s.pdf' % num,pagesize=A4)
        xlength,ylength = A4
        #print('width:%d high:%d'%(xlength,ylength))
        #c.line(1,1,ylength/2,ylength)
        #设置文字类型及字号
        #c.setFont('hei',20)
        #生成一个table表格
        atable = [[1,2,3,4,5,6,7,8],[11,12,13,14,15,16,17,18]]
        t = Table(atable,50,20)
        t.setStyle(TableStyle([('ALIGN',(0,0),(3,1),'CENTER'),
                               ('INNERGRID',(0,0),(-1,-1),0.25,colors.black),
                               ('BOX',(0,0),(-1,-1),0.25,colors.black)]))
        textOb = c.beginText(1,ylength-10)
        indexVlaue = 0
        while(indexVlaue < ylength):
            textStr = u'''wo shi tupian---wo shi tupian--wo shi tupian--wo shi tupian%d'''%indexVlaue + salt
            #print('nextline,nextline%d'%indexVlaue)
            textOb.textLine(textStr)
            indexVlaue = indexVlaue + 1
            break
        c.drawText(textOb)
        
        #简单的图片载入
        imageValue = 'file\\dlp.png'
        c.drawImage(imageValue,97,97,650,650)
        #c.drawImage('file\\dlp.png',50,50,50,50)
        t.split(0,0)
        t.drawOn(c,100,1)
        c.showPage()
        #换页的方式不同的showPage
        while x < s:
            imageValue = 'file\\dlp.png'
            c.drawImage(imageValue,97,97,650,650)
            c.drawString(0,0,'tupian%s' % x)
            c.showPage()
            x = x + 1
        c.save()
    def CreatePdfE(self,s,n):
        m = 0
        o = office()
        while m < n:
            o.CreatePdfNum(s,m)
            m = m + 1
        return True

    def GenSizeFile(self,fileSize):
        #file path
        filePath="Result\\dlpsize.txt"
        salt = ''.join(random.sample(string.ascii_letters + string.digits, 16))
        
        # 生成固定大小的文件
        # date size
        #ds = 0
        with open(filePath, "w") as f:
            '''
            while ds < fileSize:
                f.write(str(round(random.uniform(-1000, 1000),2)))
                f.write("\n")
                ds = os.path.getsize(filePath)
            print ds
            '''
            f.write('王宝强%s' % salt)
            f.seek(1024*1024*fileSize)
            f.write('\x00')
            f.write('\n')
            #f.write('王宝强')
            f.close()
        return True
        # print(os.path.getsize(filePath))

    # 支持文件类型
    # 用16进制字符串的目的是可以知道文件头是多少字节
    # 各种文件头的长度不一样，少半2字符，长则8字符
    def typeList(self):
        return {
            "574C4653": WLFS
            }
 
    # 字节码转16进制字符串
    def bytes2hex(self,bytes):
        num = len(bytes)
        hexstr = u""
        for i in range(num):
            t = u"%x" % bytes[i]
            if len(t) % 2:
                hexstr += u"0"
            hexstr += t
        return hexstr.upper()
    '''

    # 字节码转16进制字符串
    def bytes2hex(self,bytes):
        print(u'关键码转码……')
        num = len(bytes)
        hexstr = u""
        for i in range(num):
            t = u"%x" % bytes[i]
            if len(t) % 2:
                hexstr += u"0"
            hexstr += t
        return hexstr.upper()
    '''
    def get_desktop(self):
        key = _winreg.OpenKey(_winreg.HKEY_CURRENT_USER,r'Software\Microsoft\Windows\CurrentVersion\Explorer\Shell Folders')
        return _winreg.QueryValueEx(key, "Desktop")[0]
    '''
    # 获取文件类型
    def filetype(self,filename):
        app = office()
        print(u'读文件二进制码中……')
        binfile = open(filename, 'rb') # 必需二制字读取
        print(u'提取关键码……')
        bins = binfile.read(20) #提取20个字符
        binfile.close() #关闭文件流
        bins = app.bytes2hex(bins) #转码
        bins = bins.lower()#小写
        print(bins)
        tl = app.typeList()#文件类型
        ftype = 'unknown'
        print(u'关键码比对中……')
        for hcode in tl.keys():
            lens = len(hcode) # 需要的长度
            if bins[0:lens] == hcode:
                ftype = tl[hcode]
                break
        if ftype == 'unknown':#全码未找到，优化处理，码表取5位验证
            bins = bins[0:5]
            for hcode in tl.keys():
                if len(hcode) > 5 and bins == hcode[0:5]:
                    ftype = tl[hcode]
                    break
        return ftype
    ''' 
    # 获取文件类型
    def filetype(self,filename):
        app = office()
        binfile = open(filename, 'rb') # 必需二制字读取
        tl = app.typeList()
        ftype = 'unknown'
        for hcode in tl.keys():
            numOfBytes = len(hcode) / 2 # 需要读多少字节
            binfile.seek(0) # 每次读取都要回到文件头，不然会一直往后读取
            hbytes = struct.unpack_from("B"*numOfBytes, binfile.read(numOfBytes)) # 一个 "B"表示一个字节
            f_hcode = app.bytes2hex(hbytes)
            if f_hcode == hcode:
                ftype = tl[hcode]
                break
        binfile.close()
        return ftype
    
#end--------------------------------------------------------
if __name__ == "__main__":
    app = office()
    def usage():
        print(u'\
        -h or --help：显示帮助信息\n\
        -e ：创建Excel文件\n\
        -w ：创建Word文件\n\
        -p ：创建Ppt文件\n\
        -t ：创建Txt文件\n\
        -i ：创建Png文件\n\
        -P ：创建Pdf文件\n\
        -v or --version：显示版本\
        ')
    if len(sys.argv) == 1:
        usage()
        sys.exit()

    try:
        opts, args = getopt.getopt(sys.argv[1:], "hewptzifPSv", ["help", "version"])
    except getopt.GetoptError:
        usage()        

    for cmd, arg in opts:
        if cmd in ("-h", "--help"):
            usage()
        elif cmd in ("-e", "1"):
            text = "我不是关键字,我是来占格子的"
            s = 100
            app.CreateExcel(text,s)
        elif cmd in ("-w", "1"):
            s = 100
            app.CreateWord(s)
        elif cmd in ("-p", "1"):
            s = 100
            app.CreatePpt(s)
        elif cmd in ("-t", "1"):
            s = 500
            app.CreateTxt(s)
        elif cmd in ("-z", "1"):
            text = "我不是关键字,我是来占格子的"
            s = 10
            n = 10
            app.CreatePdfE(s,n)
        elif cmd in ("-i", "1"):
            app.CreateImage()
        elif cmd in ("-f", "1"):
            path = app.get_desktop()
            filename = path + '\dlp.txt'
            print app.filetype(filename)
        elif cmd in ("-P", "1"):
            s = 10
            app.CreatePdf(s)
        elif cmd in ("-S", "1"):
            app.GenSizeFile(1)
        elif cmd in ("-v", "--version"):
            print("version 1.0")